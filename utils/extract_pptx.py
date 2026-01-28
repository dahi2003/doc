from pptx import Presentation

def extract_text_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text.append(shape.text_frame.text)
    return " ".join(text)